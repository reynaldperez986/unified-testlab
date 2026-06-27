from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("ai_agent_flow_search_download_delete_update.pptx")
EXEC_OUTPUT_PATH = Path("ai_agent_flow_executive_summary.pptx")
CLIENT_OUTPUT_PATH = Path("ai_agent_flow_client_facing.pptx")
ASSET_DIR = Path("generated_assets")
ICON_DIR = ASSET_DIR / "ppt_icons"
SCREENSHOT_DIR = Path("logs") / "screenshots"
UI_SCREENSHOTS = [
    SCREENSHOT_DIR / "aacb4887-38cd-48e9-aacb-34cbd944fdf1_step1_pass.png",
    SCREENSHOT_DIR / "aacb4887-38cd-48e9-aacb-34cbd944fdf1_step3_pass.png",
    SCREENSHOT_DIR / "aacb4887-38cd-48e9-aacb-34cbd944fdf1_step6_pass.png",
]

BG = RGBColor(245, 247, 250)
INK = RGBColor(34, 40, 49)
MUTED = RGBColor(92, 99, 112)
ACCENT = RGBColor(17, 138, 178)
ACCENT_DARK = RGBColor(7, 59, 76)
SEARCH = RGBColor(0, 109, 119)
DOWNLOAD = RGBColor(42, 157, 143)
DELETE = RGBColor(214, 40, 40)
UPDATE = RGBColor(244, 162, 97)
CARD = RGBColor(255, 255, 255)
ARROW = RGBColor(150, 159, 171)
SOFT = RGBColor(227, 233, 240)


def _rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return int(color[0]), int(color[1]), int(color[2])


def _hex_color(color: RGBColor) -> str:
    return "#%02x%02x%02x" % _rgb_tuple(color)


def _ensure_visual_assets() -> dict[str, Path]:
    from PIL import Image, ImageDraw, ImageFont

    ICON_DIR.mkdir(parents=True, exist_ok=True)
    specs = {
        "request": {"symbol": "C", "color": ACCENT_DARK, "accent": ACCENT, "label": "Request"},
        "routing": {"symbol": "R", "color": ACCENT, "accent": SEARCH, "label": "Routing"},
        "llm": {"symbol": "AI", "color": SEARCH, "accent": ACCENT_DARK, "label": "Assist"},
        "action": {"symbol": "A", "color": DOWNLOAD, "accent": SEARCH, "label": "Action"},
        "reply": {"symbol": "OK", "color": DELETE, "accent": UPDATE, "label": "Reply"},
        "fast": {"symbol": "F", "color": ACCENT, "accent": ACCENT_DARK, "label": "Fast"},
        "safe": {"symbol": "S", "color": DELETE, "accent": ACCENT_DARK, "label": "Safe"},
        "clear": {"symbol": "CL", "color": SEARCH, "accent": DOWNLOAD, "label": "Clear"},
        "consistent": {"symbol": "CN", "color": DOWNLOAD, "accent": ACCENT, "label": "Consistent"},
    }

    def _load_font(size: int, bold: bool = False):
        candidates = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
            "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
        ]
        for candidate in candidates:
            if Path(candidate).exists():
                return ImageFont.truetype(candidate, size=size)
        return ImageFont.load_default()

    assets: dict[str, Path] = {}
    for key, spec in specs.items():
        png_path = ICON_DIR / f"{key}.png"
        svg_path = ICON_DIR / f"{key}.svg"

        image = Image.new("RGBA", (256, 256), (245, 247, 250, 0))
        draw = ImageDraw.Draw(image)
        card_color = _rgb_tuple(CARD)
        main_color = _rgb_tuple(spec["color"])
        accent_color = _rgb_tuple(spec["accent"])
        soft_color = _rgb_tuple(SOFT)
        ink_color = _rgb_tuple(INK)

        draw.rounded_rectangle((10, 10, 246, 246), radius=36, fill=card_color, outline=soft_color, width=4)
        draw.rounded_rectangle((24, 24, 232, 232), radius=30, fill=(255, 255, 255, 255), outline=soft_color, width=2)
        draw.ellipse((46, 42, 210, 206), fill=soft_color)
        draw.ellipse((62, 58, 194, 190), fill=main_color)
        draw.ellipse((42, 40, 76, 74), fill=accent_color)
        draw.rounded_rectangle((48, 194, 208, 222), radius=10, fill=accent_color)

        symbol_font = _load_font(56, bold=True)
        label_font = _load_font(22, bold=True)
        bbox = draw.textbbox((0, 0), spec["symbol"], font=symbol_font)
        sw = bbox[2] - bbox[0]
        sh = bbox[3] - bbox[1]
        draw.text((128 - sw / 2, 122 - sh / 2), spec["symbol"], fill=(255, 255, 255, 255), font=symbol_font)
        label_bbox = draw.textbbox((0, 0), spec["label"], font=label_font)
        lw = label_bbox[2] - label_bbox[0]
        lh = label_bbox[3] - label_bbox[1]
        draw.text((128 - lw / 2, 207 - lh / 2), spec["label"], fill=(255, 255, 255, 255), font=label_font)
        image.save(png_path)

        svg_path.write_text(
            """
<svg xmlns="http://www.w3.org/2000/svg" width="256" height="256" viewBox="0 0 256 256">
  <rect x="10" y="10" width="236" height="236" rx="36" fill="{card}" stroke="{soft}" stroke-width="4"/>
  <rect x="24" y="24" width="208" height="208" rx="30" fill="#ffffff" stroke="{soft}" stroke-width="2"/>
  <circle cx="128" cy="124" r="82" fill="{soft}"/>
  <circle cx="128" cy="124" r="66" fill="{main}"/>
  <circle cx="59" cy="57" r="17" fill="{accent}"/>
  <rect x="48" y="194" width="160" height="28" rx="10" fill="{accent}"/>
  <text x="128" y="135" text-anchor="middle" font-family="Arial, Segoe UI, sans-serif" font-size="56" font-weight="700" fill="#ffffff">{symbol}</text>
  <text x="128" y="214" text-anchor="middle" font-family="Arial, Segoe UI, sans-serif" font-size="22" font-weight="700" fill="#ffffff">{label}</text>
</svg>
            """.strip().format(
                card=_hex_color(CARD),
                soft=_hex_color(SOFT),
                main=_hex_color(spec["color"]),
                accent=_hex_color(spec["accent"]),
                symbol=spec["symbol"],
                label=spec["label"],
            ),
            encoding="utf-8",
        )
        assets[key] = png_path
    return assets


def _save_with_fallback(prs, output_path: Path) -> Path:
    try:
        prs.save(output_path)
        return output_path.resolve()
    except PermissionError:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        fallback = output_path.with_stem(f"{output_path.stem}_{stamp}")
        prs.save(fallback)
        return fallback.resolve()


def _set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_title(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.65))
    text_frame = title_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = ACCENT_DARK

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.92), Inches(11.5), Inches(0.45))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        paragraph = subtitle_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED


def _add_card(slide, left, top, width, height, title, body_lines, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    title_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.45))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = color
    title_bar.line.color.rgb = color

    tf = title_bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    body = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.58), width - Inches(0.3), height - Inches(0.72))
    frame = body.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.clear()

    for index, line in enumerate(body_lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(6)


def _add_chevron(slide, left, top, width, height, color=ARROW):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def _add_text_block(slide, left, top, width, height, lines, *, font_name="Aptos", font_size=11, color=INK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
    return box


def _add_icon_plate(slide, left, top, size, symbol, title, color, *, subtitle="", symbol_font="Segoe UI Symbol", icon_key=""):
    plate = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, size, size)
    plate.fill.solid()
    plate.fill.fore_color.rgb = CARD
    plate.line.color.rgb = SOFT
    plate.line.width = Pt(1.25)

    if icon_key:
        assets = _ensure_visual_assets()
        icon_path = assets.get(icon_key)
        if icon_path and icon_path.exists():
            slide.shapes.add_picture(str(icon_path), left + Inches(0.08), top + Inches(0.08), width=size - Inches(0.16), height=size - Inches(0.16))
    else:
        halo = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.12), top + Inches(0.12), size - Inches(0.24), size - Inches(0.24))
        halo.fill.solid()
        halo.fill.fore_color.rgb = SOFT
        halo.line.color.rgb = SOFT

        core = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.28), top + Inches(0.28), size - Inches(0.56), size - Inches(0.56))
        core.fill.solid()
        core.fill.fore_color.rgb = color
        core.line.color.rgb = color

        spark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.08), top + Inches(0.08), Inches(0.18), Inches(0.18))
        spark.fill.solid()
        spark.fill.fore_color.rgb = color
        spark.line.color.rgb = color

        symbol_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), size - Inches(0.4), size - Inches(0.28))
        tf = symbol_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = 1
        r = p.add_run()
        r.text = symbol
        r.font.name = symbol_font
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)

    title_box = slide.shapes.add_textbox(left - Inches(0.05), top + size + Inches(0.08), size + Inches(0.1), Inches(0.3))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = 1
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = INK

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(left - Inches(0.15), top + size + Inches(0.34), size + Inches(0.3), Inches(0.34))
        tf = subtitle_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = 1
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(8)
        r.font.color.rgb = MUTED

    return plate


def _add_flow_connector(slide, left, top, width, height, color=ARROW, label=""):
    connector = _add_chevron(slide, left, top, width, height, color)
    if label:
        _add_text_block(slide, left - Inches(0.15), top - Inches(0.28), width + Inches(0.3), Inches(0.22), [label], font_size=8, color=MUTED)
    return connector


def _add_mapping_card(slide, left, top, width, height, title, lines, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = color
    shape.line.width = Pt(1.25)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.38))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.color.rgb = color

    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    _add_text_block(
        slide,
        left + Inches(0.14),
        top + Inches(0.48),
        width - Inches(0.28),
        height - Inches(0.6),
        lines,
        font_name="Consolas",
        font_size=8,
        color=INK,
    )


def _add_screenshot_panel(slide, left, top, width, height, image_path: Path, title: str, caption: str):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD
    frame.line.color.rgb = SOFT
    frame.line.width = Pt(1.25)

    image_top = top + Inches(0.16)
    image_height = height - Inches(0.9)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), left + Inches(0.12), image_top, width=width - Inches(0.24), height=image_height)

    _add_text_block(slide, left + Inches(0.15), top + height - Inches(0.58), width - Inches(0.3), Inches(0.2), [title], font_size=10, color=ACCENT_DARK, bold=True)
    _add_text_block(slide, left + Inches(0.15), top + height - Inches(0.34), width - Inches(0.3), Inches(0.24), [caption], font_size=8, color=MUTED)


def _add_process_strip(slide, steps, left, top, total_width, height, color):
    step_width = total_width / len(steps)
    for index, label in enumerate(steps):
        box_left = left + step_width * index
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            box_left,
            top,
            step_width - Inches(0.08),
            height,
        )
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = color
        box.line.width = Pt(1.25)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = INK
        if index < len(steps) - 1:
            _add_chevron(
                slide,
                box_left + step_width - Inches(0.28),
                top + Inches(0.18),
                Inches(0.22),
                Inches(0.34),
                color,
            )


def _add_flow_slide(prs, title, subtitle, color, steps, examples):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, title, subtitle)

    left = Inches(0.65)
    top = Inches(1.45)
    card_width = Inches(2.8)
    card_height = Inches(1.65)
    gap = Inches(0.18)

    for index, (step_title, lines) in enumerate(steps):
        _add_card(
            slide,
            left + index * (card_width + gap),
            top,
            card_width,
            card_height,
            step_title,
            lines,
            color,
        )
        if index < len(steps) - 1:
            _add_chevron(
                slide,
                left + (index + 1) * card_width + index * gap + Inches(0.02),
                top + Inches(0.57),
                Inches(0.12),
                Inches(0.34),
                color,
            )

    _add_process_strip(
        slide,
        ["User asks", "Parser checks", "LLM helps if needed", "Tool executes", "Reply rendered"],
        Inches(0.65),
        Inches(3.0),
        Inches(11.55),
        Inches(0.42),
        color,
    )

    example_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(3.62),
        Inches(11.55),
        Inches(2.76),
    )
    example_box.fill.solid()
    example_box.fill.fore_color.rgb = CARD
    example_box.line.color.rgb = color
    example_box.line.width = Pt(1.5)

    tf = example_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Sample Phrases"
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = color

    for phrase in examples:
        p = tf.add_paragraph()
        p.text = phrase
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = INK
        p.bullet = True

    return slide


def _add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Architecture View", "The agent uses deterministic routing first, then selectively asks the LLM to help with ambiguous intent and response rendering.")

    _add_card(slide, Inches(0.7), Inches(1.65), Inches(2.25), Inches(1.5), "User Input", ["Natural language request", "Optional conversation history"], ACCENT_DARK)
    _add_card(slide, Inches(3.1), Inches(1.65), Inches(2.45), Inches(1.5), "Deterministic Layer", ["Regex parsers", "Fast trigger tuples", "Direct routing for obvious phrasing"], ACCENT)
    _add_card(slide, Inches(5.8), Inches(1.3), Inches(2.45), Inches(2.2), "Guarded LLM Layer", ["Boolean trigger classifier", "Fallback intent parser", "Only called when direct parsing misses"], SEARCH)
    _add_card(slide, Inches(8.55), Inches(1.65), Inches(2.1), Inches(1.5), "Tool Layer", ["Search", "Download", "Delete", "Update"], DOWNLOAD)
    _add_card(slide, Inches(10.9), Inches(1.65), Inches(1.75), Inches(1.5), "Finalizer", ["LLM rewrites raw tool output into user-facing text"], DELETE)

    _add_chevron(slide, Inches(2.97), Inches(2.2), Inches(0.18), Inches(0.32), ACCENT)
    _add_chevron(slide, Inches(5.56), Inches(2.2), Inches(0.18), Inches(0.32), SEARCH)
    _add_chevron(slide, Inches(8.3), Inches(2.2), Inches(0.18), Inches(0.32), DOWNLOAD)
    _add_chevron(slide, Inches(10.67), Inches(2.2), Inches(0.18), Inches(0.32), DELETE)

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.4), Inches(1.55))
    note.fill.solid()
    note.fill.fore_color.rgb = CARD
    note.line.color.rgb = SOFT
    tf = note.text_frame
    tf.clear()
    points = [
        "Common phrasing never needs the LLM at the trigger stage.",
        "Ambiguous phrasing is screened by a small yes/no classifier before the heavier fallback parser runs.",
        "Tool execution remains deterministic against PostgreSQL and generated downloads.",
        "The final response LLM improves readability without changing the underlying action result.",
    ]
    for index, point in enumerate(points):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = point
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = INK
        p.bullet = True


def _add_hybrid_identity_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(
        slide,
        "Hybrid Agentic Generative AI",
        "ai_agent.py is a hybrid agentic generative AI assistant with deterministic routing plus LLM-based fallback and tool execution.",
    )

    stages = [
        ("◔", "User Input", "Natural-language request enters the assistant.", ACCENT_DARK, "Request"),
        ("◇", "Deterministic Routing", "Regex and rule-based parsers handle the common paths first.", ACCENT, "Rules"),
        ("☁", "LLM Fallback", "A constrained LLM classifier and parser help only when the deterministic layer misses.", SEARCH, "Assist"),
        ("⚙", "Tool Execution", "Search, download, delete, update, and project actions run against the system.", DOWNLOAD, "Action"),
        ("✦", "Final Response", "The LLM rewrites the raw tool output into a concise user-facing answer.", DELETE, "Reply"),
    ]

    start_left = Inches(0.63)
    top = Inches(1.72)
    icon_size = Inches(0.95)
    card_width = Inches(2.04)
    card_height = Inches(2.2)
    gap = Inches(0.4)

    for index, (icon, title, body, color, subtitle) in enumerate(stages):
        left = start_left + index * (card_width + gap)
        _add_icon_plate(slide, left + Inches(0.48), top, icon_size, icon, title, color, subtitle=subtitle, icon_key=("request", "routing", "llm", "action", "reply")[index])
        _add_card(
            slide,
            left,
            top + Inches(1.38),
            card_width,
            card_height,
            title,
            [body],
            color,
        )
        if index < len(stages) - 1:
            _add_flow_connector(
                slide,
                left + card_width + Inches(0.05),
                top + Inches(1.75),
                Inches(0.28),
                Inches(0.4),
                ARROW,
                label="next",
            )

    summary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(5.35),
        Inches(11.45),
        Inches(1.15),
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = CARD
    summary.line.color.rgb = SOFT
    tf = summary.text_frame
    tf.clear()
    bullets = [
        "Generative AI: the LLM interprets intent, generates constrained JSON, and rewrites final responses.",
        "Agentic AI: the assistant can choose and execute real tools instead of only answering with text.",
        "Deterministic first: common commands stay fast, predictable, and cheaper than an LLM-first design.",
    ]
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = INK
        p.bullet = True


def _add_function_mapping_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(
        slide,
        "Function Mapping",
        "Exact ai_agent.py functions mapped to the flowchart blocks used by the hybrid assistant.",
    )

    mapping_cards = [
        (
            Inches(0.65), Inches(1.45), Inches(3.95), Inches(1.7), "1. Entry + Orchestration", [
                "handle_chat_message",
                "_execute_tool_interaction",
                "TOOL_DISPATCH",
            ], ACCENT_DARK,
        ),
        (
            Inches(4.7), Inches(1.45), Inches(3.95), Inches(1.7), "2. Deterministic Parsers", [
                "_parse_delete_request, _is_delete_confirmation",
                "_pending_delete_from_history",
                "_parse_show_steps_request",
                "_parse_download_request",
                "_parse_update_step_request",
                "_parse_bulk_update_data_request",
                "_parse_update_data_value_request",
                "_parse_search_sessions_request",
            ], ACCENT,
        ),
        (
            Inches(8.75), Inches(1.45), Inches(3.95), Inches(1.7), "3. LLM Trigger Gates", [
                "_looks_like_delete_candidate",
                "_looks_like_show_steps_candidate",
                "_looks_like_download_candidate",
                "_looks_like_update_candidate",
                "_looks_like_search_candidate",
                "_llm_should_try_*_parse",
                "_should_try_llm_*_parse",
            ], SEARCH,
        ),
        (
            Inches(0.65), Inches(3.45), Inches(3.95), Inches(2.0), "4. LLM Fallback Parsers", [
                "_llm_parse_delete_request",
                "_llm_parse_show_steps_request",
                "_llm_parse_download_request",
                "_llm_parse_update_request",
                "_llm_parse_search_sessions_request",
                "_call_ollama, _extract_json_object",
            ], DOWNLOAD,
        ),
        (
            Inches(4.7), Inches(3.45), Inches(3.95), Inches(2.0), "5. Tool Actions", [
                "tool_delete_test_case",
                "tool_show_steps",
                "tool_download_session",
                "tool_update_step",
                "tool_bulk_update_data",
                "tool_update_data_value",
                "tool_search_sessions",
            ], DELETE,
        ),
        (
            Inches(8.75), Inches(3.45), Inches(3.95), Inches(2.0), "6. Response Rendering", [
                "_tool_result",
                "_render_tool_result_with_llm",
                "_render_direct_reply_with_llm",
                "_clean_final_reply",
                "_build_history_text",
                "_extract_tool_call",
            ], UPDATE,
        ),
    ]

    for left, top, width, height, title, lines, color in mapping_cards:
        _add_mapping_card(slide, left, top, width, height, title, lines, color)

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.8), Inches(11.8), Inches(0.75))
    footer.fill.solid()
    footer.fill.fore_color.rgb = CARD
    footer.line.color.rgb = SOFT
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "The detailed deck shows the exact implementation touchpoints; the client-facing deck removes these code references and focuses on service outcomes, governance, and user experience."
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def _add_end_to_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "End-to-End Flowchart", "One operating pattern handles search, download, delete, and update from request intake through user-ready response.")

    stages = [
        (Inches(0.55), "Request Intake", "◔", ACCENT_DARK, ["User asks in natural language", "History keeps follow-up context"]),
        (Inches(3.0), "Intent Identification", "◇", ACCENT, ["Deterministic parser first", "LLM gate only for ambiguous phrasing"]),
        (Inches(5.45), "Business Action", "⚙", SEARCH, ["Search, download, delete, or update", "Database and file operations stay deterministic"]),
        (Inches(7.9), "Structured Result", "▣", DOWNLOAD, ["Status, tables, suggestions, download metadata", "Delete stays confirmation-first"]),
        (Inches(10.35), "Clear Response", "✦", DELETE, ["LLM rewrites output for readability", "Business intent is preserved"]),
    ]

    for index, (left, title, symbol, color, lines) in enumerate(stages):
        _add_icon_plate(slide, left + Inches(0.52), Inches(1.45), Inches(0.88), symbol, title, color, icon_key=("request", "routing", "action", "reply", "reply")[index])
        _add_card(slide, left, Inches(2.35), Inches(2.05), Inches(1.92), title, lines, color)
        if index < len(stages) - 1:
            _add_flow_connector(slide, left + Inches(2.08), Inches(3.02), Inches(0.28), Inches(0.34), ARROW, label="shared runtime")

    lane_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.75), Inches(12.0), Inches(1.45))
    lane_box.fill.solid()
    lane_box.fill.fore_color.rgb = CARD
    lane_box.line.color.rgb = SOFT
    lane_box.line.width = Pt(1.25)

    _add_text_block(slide, Inches(0.95), Inches(4.98), Inches(1.3), Inches(0.25), ["Operation lanes"], font_size=11, color=ACCENT_DARK, bold=True)
    lanes = [
        (SEARCH, "Search", "Find the right test case quickly"),
        (DOWNLOAD, "Download", "Deliver the script in the requested format"),
        (DELETE, "Delete", "Require review before permanent removal"),
        (UPDATE, "Update", "Apply controlled changes with before/after feedback"),
    ]
    for index, (color, label, message) in enumerate(lanes):
        left = Inches(2.2 + index * 2.55)
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(4.95), Inches(2.25), Inches(0.34))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.color.rgb = color
        tf = pill.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = 1
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        _add_text_block(slide, left, Inches(5.38), Inches(2.25), Inches(0.45), [message], font_size=8, color=MUTED)


def _add_executive_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.8), Inches(11.95), Inches(1.4))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_DARK
    banner.line.color.rgb = ACCENT_DARK
    tf = banner.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "AI Agent Executive Summary"
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "Search, Download, Delete, and Update | WebConX Automation Platform"
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230, 236, 240)

    _add_process_strip(
        slide,
        ["User request", "Deterministic route", "LLM assist if ambiguous", "Tool action", "Final reply"],
        Inches(0.9),
        Inches(2.75),
        Inches(11.4),
        Inches(0.55),
        ACCENT,
    )

    _add_card(slide, Inches(0.95), Inches(3.7), Inches(3.5), Inches(2.0), "What Changed", ["Intent handling is now hybrid instead of purely prompt-driven.", "Each major action has deterministic routing plus a guarded LLM fallback."], ACCENT)
    _add_card(slide, Inches(4.8), Inches(3.7), Inches(3.5), Inches(2.0), "Business Value", ["Faster common-path responses", "Lower LLM usage", "Better handling of natural language variation"], SEARCH)
    _add_card(slide, Inches(8.65), Inches(3.7), Inches(3.5), Inches(2.0), "Operational Value", ["DB actions stay deterministic.", "Ambiguous requests still route correctly.", "Final replies stay user-friendly."], DOWNLOAD)


def _add_executive_operations_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Operational View", "Each operation now follows the same architecture, with a different parser and tool at the end.")

    items = [
        (SEARCH, "Search", ["Explicit parser first", "LLM trigger for existence-style questions", "Search parser extracts query", "DB returns matching scripts"]),
        (DOWNLOAD, "Download", ["Parser resolves name, folder, format", "LLM trigger catches send/share phrasing", "Tool handles duplicate names", "Returns URL or asks for format/project"]),
        (DELETE, "Delete", ["Confirmation-first workflow", "LLM trigger catches archive/drop phrasing", "Tool shows summary before deletion", "Second message confirms actual removal"]),
        (UPDATE, "Update", ["Single, bulk, and step updates parsed separately", "LLM trigger catches revise/fix wording", "Folder context preserved", "Before/after result returned"]),
    ]

    for index, (color, title, lines) in enumerate(items):
        row = index // 2
        col = index % 2
        _add_card(slide, Inches(0.8 + col * 6.2), Inches(1.55 + row * 2.35), Inches(5.5), Inches(1.95), title, lines, color)


def _add_executive_takeaways_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Key Takeaways", "The platform now uses the LLM as a selective assistant instead of the only router.")

    _add_card(slide, Inches(0.9), Inches(1.6), Inches(3.65), Inches(3.8), "Decision Logic", ["Deterministic first", "LLM only for ambiguous trigger detection and fallback parsing", "Final rewrite improves clarity"], ACCENT_DARK)
    _add_card(slide, Inches(4.85), Inches(1.6), Inches(3.65), Inches(3.8), "Risk Control", ["Lower accidental cross-intent matches", "Confirmation guard for delete", "Folder/project context preserved for duplicate names"], DELETE)
    _add_card(slide, Inches(8.8), Inches(1.6), Inches(3.65), Inches(3.8), "Executive Outcome", ["More reliable automation assistant", "Lower model cost on common paths", "Better readiness for broader natural-language use"], SEARCH)


def _add_client_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.72), Inches(11.95), Inches(1.45))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_DARK
    banner.line.color.rgb = ACCENT_DARK
    tf = banner.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "WebConX AI Executive Brief"
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "Executive audience | Faster service resolution, safer operations, clearer client experience"
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230, 236, 240)

    _add_card(slide, Inches(0.95), Inches(2.55), Inches(3.55), Inches(3.2), "Executive View", ["The assistant shortens turnaround for routine automation requests.", "Users receive clearer next steps when information is missing.", "High-risk actions keep explicit human confirmation."], ACCENT)
    _add_card(slide, Inches(4.9), Inches(2.55), Inches(3.55), Inches(3.2), "Operational Outcome", ["Search, export, update, and retirement workflows use one consistent service pattern.", "The platform handles duplicate names and folder context more reliably.", "Responses stay business-friendly even when requests are informal."], SEARCH)
    _add_card(slide, Inches(8.85), Inches(2.55), Inches(3.05), Inches(3.2), "Why It Matters", ["Lower support effort", "Better client confidence", "More predictable service quality"], DOWNLOAD)


def _add_client_value_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Executive Framing", "The experience is designed to improve service speed and governance together, not trade one for the other.")

    _add_icon_plate(slide, Inches(1.0), Inches(1.7), Inches(1.0), "⚡", "Fast", ACCENT, subtitle="Common requests", icon_key="fast")
    _add_icon_plate(slide, Inches(4.25), Inches(1.7), Inches(1.0), "✓", "Safe", DELETE, subtitle="Controlled actions", icon_key="safe")
    _add_icon_plate(slide, Inches(7.5), Inches(1.7), Inches(1.0), "✦", "Clear", SEARCH, subtitle="Readable replies", icon_key="clear")
    _add_icon_plate(slide, Inches(10.75), Inches(1.7), Inches(1.0), "↺", "Consistent", DOWNLOAD, subtitle="Shared runtime", icon_key="consistent")

    _add_card(slide, Inches(0.8), Inches(3.2), Inches(3.0), Inches(2.1), "Faster resolution", ["Routine requests move faster because the platform recognizes the common patterns immediately.", "That improves response time without adding more manual triage."], ACCENT)
    _add_card(slide, Inches(3.95), Inches(3.2), Inches(3.0), Inches(2.1), "Stronger control", ["Sensitive actions such as deletion remain governed by confirmation and application rules.", "The AI improves guidance, not authority."], DELETE)
    _add_card(slide, Inches(7.1), Inches(3.2), Inches(3.0), Inches(2.1), "Better client experience", ["People can ask naturally and still get a clear, business-readable answer.", "That lowers friction for self-service usage."], SEARCH)
    _add_card(slide, Inches(10.25), Inches(3.2), Inches(2.25), Inches(2.1), "Scalable model", ["One service pattern supports multiple workflows."], DOWNLOAD)


def _add_client_operations_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Service Scenarios", "From an executive lens, the value comes from faster handling of the most common client requests.")

    items = [
        (SEARCH, "Find", ["Locate the right automation asset quickly", "Useful when stakeholders know only part of the name or project"]),
        (DOWNLOAD, "Share", ["Provide an export in the requested format", "Reduces manual handoff and speeds client follow-through"]),
        (UPDATE, "Change", ["Apply requested revisions with visible confirmation", "Supports ongoing maintenance without a heavy support loop"]),
        (DELETE, "Retire", ["Review before removing outdated assets", "Helps prevent avoidable operational loss"]),
    ]
    for index, (color, title, lines) in enumerate(items):
        row = index // 2
        col = index % 2
        _add_card(slide, Inches(0.95 + col * 5.95), Inches(1.75 + row * 2.15), Inches(5.35), Inches(1.8), title, lines, color)


def _add_client_controls_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Governance View", "The assistant strengthens front-end experience while core platform rules still govern the action itself.")

    _add_card(slide, Inches(0.9), Inches(1.65), Inches(3.7), Inches(3.6), "Control Points", ["Application logic still decides the final action.", "Deletion remains confirmation-first.", "Folder context reduces the risk of acting on the wrong asset when names repeat."], ACCENT_DARK)
    _add_card(slide, Inches(4.85), Inches(1.65), Inches(3.7), Inches(3.6), "Business Benefits", ["Less rework from misunderstood requests", "More predictable service delivery", "Higher confidence in self-service interactions"], SEARCH)
    _add_card(slide, Inches(8.8), Inches(1.65), Inches(3.7), Inches(3.6), "Executive Outcome", ["A more polished client-facing support layer", "Faster handling of common operational work", "A credible foundation for broader AI-assisted service offerings"], DOWNLOAD)


def _add_client_ui_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Product UI In Context", "Real screenshots from the automation flow make the assistant story concrete: data entry, navigation, and resulting state are all visible in the product itself.")

    panels = [
        (Inches(0.55), UI_SCREENSHOTS[0], "Capture", "The workflow begins inside the real application UI, not in a synthetic demo state."),
        (Inches(4.4), UI_SCREENSHOTS[1], "Progress", "Recorded interaction shows typed values and action execution as the script advances."),
        (Inches(8.25), UI_SCREENSHOTS[2], "Result", "The final state provides visible evidence that the underlying workflow completed successfully."),
    ]
    for left, image_path, title, caption in panels:
        _add_screenshot_panel(slide, left, Inches(1.45), Inches(3.45), Inches(4.7), image_path, title, caption)

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.62))
    footer.fill.solid()
    footer.fill.fore_color.rgb = CARD
    footer.line.color.rgb = SOFT
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "These screenshots come directly from the recorded run artifacts in logs/screenshots, so the deck ends with product evidence rather than abstract architecture alone."
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def build_executive_summary():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_executive_title_slide(prs)
    _add_hybrid_identity_slide(prs)
    _add_architecture_slide(prs)
    _add_end_to_end_slide(prs)
    _add_executive_operations_slide(prs)
    _add_executive_takeaways_slide(prs)

    return _save_with_fallback(prs, EXEC_OUTPUT_PATH)


def build_client_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_client_title_slide(prs)
    _add_client_value_slide(prs)
    _add_client_operations_slide(prs)
    _add_end_to_end_slide(prs)
    _add_client_controls_slide(prs)
    _add_client_ui_slide(prs)

    return _save_with_fallback(prs, CLIENT_OUTPUT_PATH)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.65), Inches(12.0), Inches(1.3))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_DARK
    banner.line.color.rgb = ACCENT_DARK
    tf = banner.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "AI Agent Flow for Search, Download, Delete, and Update"
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    p = tf.add_paragraph()
    p.text = "WebConX Automation Platform | Hybrid routing: deterministic parser -> guarded LLM fallback -> tool execution -> final LLM response"
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230, 236, 240)

    _add_card(
        slide,
        Inches(0.85),
        Inches(2.35),
        Inches(3.8),
        Inches(3.8),
        "How It Works Now",
        [
            "User message enters handle_chat_message(...).",
            "Fast regex and phrase parsers handle obvious intents first.",
            "Ambiguous phrasing uses a small yes/no LLM gate before the full fallback parser.",
            "Tool functions query PostgreSQL or generate downloads.",
            "A final LLM rewrites the raw tool result into the user-facing response.",
        ],
        ACCENT,
    )
    _add_card(
        slide,
        Inches(4.8),
        Inches(2.35),
        Inches(3.6),
        Inches(3.8),
        "Why This Hybrid Design",
        [
            "Keeps common requests fast.",
            "Reduces unnecessary Ollama calls.",
            "Lets natural language variants still work.",
            "Keeps DB actions deterministic after intent is identified.",
            "Improves final answer quality without changing tool outputs.",
        ],
        ACCENT,
    )
    _add_card(
        slide,
        Inches(8.55),
        Inches(2.35),
        Inches(3.45),
        Inches(3.8),
        "Core Layers",
        [
            "1. Deterministic parser",
            "2. LLM trigger classifier",
            "3. LLM fallback parser",
            "4. Tool dispatch",
            "5. Final response rendering",
        ],
        ACCENT,
    )

    _add_flow_slide(
        prs,
        "Search Flow",
        "Search requests try deterministic parsing first, then use an LLM trigger gate only for ambiguous search phrasing.",
        SEARCH,
        [
            ("1. Parse", ["Regex handles explicit search patterns.", "Examples: search for, find, file name=..."]),
            ("2. Gate", ["If regex misses, a boolean LLM trigger checks if the message is search-like.", "Avoids unnecessary parser calls."]),
            ("3. Extract", ["Search parser LLM returns only the cleaned query.", "Guard rejects step-centric queries unless they mention scripts."]),
            ("4. Execute", ["tool_search_sessions(query) runs the DB search.", "Results are formatted as a table."]),
        ],
        [
            "search for login",
            "can you check whether we have the smoke script",
            "please verify whether the address test case exists",
            "file name=test case 1",
        ],
    )

    _add_flow_slide(
        prs,
        "Download Flow",
        "Download routing resolves script name, optional folder, and format before returning a link or a follow-up prompt.",
        DOWNLOAD,
        [
            ("1. Parse", ["Deterministic parser reads record name, folder, and file format.", "Handles download, pull, get, copy, document."]),
            ("2. Gate", ["Ambiguous export wording uses the boolean LLM trigger.", "Examples: send, share, save as."]),
            ("3. Extract", ["Fallback parser LLM returns record_id, fmt, and folder.", "Format is normalized to csv, pdf, or doc."]),
            ("4. Execute", ["tool_download_session(...) resolves duplicates and builds the download URL.", "If duplicates exist, it asks for project/folder clarification."]),
        ],
        [
            "download login test case as pdf",
            "please send the smoke test script as csv",
            "get a copy of address from Project001",
            "save the login file as doc",
        ],
    )

    _add_flow_slide(
        prs,
        "Delete Flow",
        "Delete stays confirmation-first: identify the target, show the summary, then require an explicit confirmation before deletion.",
        DELETE,
        [
            ("1. Parse", ["Direct parser handles delete, remove, omit, ommit.", "LLM gate now also helps with archive, drop, purge, erase, wipe."]),
            ("2. Gate", ["If static triggers miss, a delete-specific LLM classifier decides if the parser should run.", "This catches softer phrasing more reliably."]),
            ("3. Extract", ["Fallback parser LLM returns record_name and optional folder.", "Project context is preserved when given."]),
            ("4. Execute", ["tool_delete_test_case(..., confirm=False) shows a summary first.", "Only a later confirmation message runs confirm=True."]),
        ],
        [
            "delete smoke test case",
            "archive this login script permanently",
            "drop the address test case from Project001",
            "remove this script permanently",
        ],
    )

    _add_flow_slide(
        prs,
        "Update Flow",
        "Update routes through deterministic single-value, bulk, and step-property parsers before using the LLM fallback.",
        UPDATE,
        [
            ("1. Parse", ["Tries update_step, then bulk_update_data, then update_data_value.", "Folder/project context is extracted when present."]),
            ("2. Gate", ["If regex misses, a boolean LLM trigger checks for revise, correct, fix, adjust, patch-style requests.", "Avoids broad matching."]),
            ("3. Extract", ["Fallback parser LLM chooses update_data_value, bulk_update_data, or update_step.", "Only explicit fields are kept."]),
            ("4. Execute", ["Update tools resolve record_id, step_no, field names, and folder.", "The DB is updated and a before/after summary is returned."]),
        ],
        [
            "update email in test case 1 to alice@example.com",
            "change step 2 locator in login test case to #submit",
            "update email to alice@example.com and address to Manila in test case 1 from Project001",
            "please revise the locator in test case 1",
        ],
    )

    _add_hybrid_identity_slide(prs)
    _add_architecture_slide(prs)
    _add_function_mapping_slide(prs)
    _add_end_to_end_slide(prs)

    return _save_with_fallback(prs, OUTPUT_PATH)


if __name__ == "__main__":
    full_output = build_presentation()
    exec_output = build_executive_summary()
    client_output = build_client_presentation()
    print(full_output)
    print(exec_output)
    print(client_output)